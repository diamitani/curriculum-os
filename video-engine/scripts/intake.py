#!/usr/bin/env python3
"""
monarch-video :: intake

Parses any source document into plain text that can be turned into a
storyboard. Handles .pptx, .pdf, .docx, .md, .txt, .csv, .json.

Usage:
    python3 intake.py FILE [FILE ...] [--out brief_source.txt]

Prints the extracted text to stdout (and optionally writes it to --out).
Missing optional parsers degrade gracefully: the file is reported as
unreadable rather than crashing the run.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

MAX_CHARS_PER_FILE = 40_000


def _pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        bits = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                bits.append(shape.text_frame.text.strip())
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        bits.append(" | ".join(cells))
        if bits:
            out.append(f"--- Slide {i} ---\n" + "\n".join(bits))
    return "\n\n".join(out)


def _pdf(path: Path) -> str:
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            return "\n\n".join(
                f"--- Page {i} ---\n{(p.extract_text() or '').strip()}"
                for i, p in enumerate(pdf.pages, 1)
                if (p.extract_text() or "").strip()
            )
    except ImportError:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n\n".join(
            f"--- Page {i} ---\n{(pg.extract_text() or '').strip()}"
            for i, pg in enumerate(reader.pages, 1)
            if (pg.extract_text() or "").strip()
        )


def _docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _json(path: Path) -> str:
    return json.dumps(json.loads(_plain(path)), indent=2)[:MAX_CHARS_PER_FILE]


PARSERS = {
    ".pptx": _pptx,
    ".pdf": _pdf,
    ".docx": _docx,
    ".md": _plain,
    ".txt": _plain,
    ".csv": _plain,
    ".json": _json,
    ".html": _plain,
}


def extract(path: Path) -> str:
    parser = PARSERS.get(path.suffix.lower())
    if parser is None:
        return f"[unsupported file type: {path.suffix} — skipped]"
    try:
        text = parser(path).strip()
    except Exception as exc:  # noqa: BLE001 — report, never crash the pipeline
        return f"[could not parse {path.name}: {type(exc).__name__}: {exc}]"
    if not text:
        return f"[{path.name} parsed but contained no extractable text]"
    if len(text) > MAX_CHARS_PER_FILE:
        text = text[:MAX_CHARS_PER_FILE] + "\n[... truncated ...]"
    return text


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract text from source documents.")
    ap.add_argument("files", nargs="+", type=Path)
    ap.add_argument("--out", type=Path, default=None)
    args = ap.parse_args()

    chunks = []
    for f in args.files:
        if not f.exists():
            chunks.append(f"===== {f} =====\n[file not found]")
            continue
        chunks.append(f"===== {f.name} =====\n{extract(f)}")

    result = "\n\n".join(chunks)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(result, encoding="utf-8")
        print(f"Wrote {len(result):,} chars → {args.out}", file=sys.stderr)
    print(result)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
